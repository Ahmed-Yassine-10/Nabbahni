# -*- coding: utf-8 -*-
"""
Translate pitch-deck-robot.html -> SentinelleRx.pptx

Content is PARSED from the HTML (never retyped), the layout maps 1:1
(1280x720 css px == 13.333x7.5 in, 1px == 9525 EMU), each stage visual is
redrawn with native shapes, and PowerPoint animations reproduce the deck's
motion: morph transitions, letter-by-letter speech, staggered card reveals.
"""

import copy
import os
import re

from bs4 import BeautifulSoup
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
PROJ = r"c:\Users\SP\OneDrive\Desktop\Pharma Suply chain"
SRC = os.path.join(PROJ, "pitch-deck-robot.html")
OUT = os.environ.get("PPTX_OUT") or os.path.join(PROJ, "SentinelleRx.pptx")

PX = 9525                      # EMU per css px  -> 1280px == 13.333in exactly
def px(v): return int(round(v * PX))

SW, SH = 1280, 720

# ---- palette lifted from the html :root ----------------------------
INK      = RGBColor(0x06, 0x10, 0x0A)
PANEL    = RGBColor(0x0C, 0x1A, 0x12)
PANEL2   = RGBColor(0x12, 0x24, 0x1A)
LINE     = RGBColor(0x1F, 0x3A, 0x2B)
FOG      = RGBColor(0x85, 0xA3, 0x92)
TEXT     = RGBColor(0xE2, 0xEF, 0xE6)
TEXT_DIM = RGBColor(0xA7, 0xC2, 0xB1)
SIGNAL   = RGBColor(0x35, 0xC7, 0x7F)
SIGNAL2  = RGBColor(0x7F, 0xF0, 0xB4)
OK       = RGBColor(0x4C, 0xAF, 0x6E)
WARN     = RGBColor(0xD9, 0xA6, 0x2E)
HOT      = RGBColor(0xDD, 0x7A, 0x38)
CRIT     = RGBColor(0xD5, 0x45, 0x5C)
PILL_W   = RGBColor(0xEA, 0xF6, 0xEE)

MONO, SANS = "Consolas", "Segoe UI"

# =====================================================================
# 1 · PARSE THE HTML
# =====================================================================
soup = BeautifulSoup(open(SRC, encoding="utf-8").read(), "html.parser")

def txt(el):
    return re.sub(r"\s+", " ", el.get_text(" ", strip=True)).strip() if el else ""

SLIDES = []
for sec in soup.select("section.slide"):
    stage = sec.select_one(".stage")
    kind = "none"
    if stage:
        for k, sel in [("radar", ".radar"), ("shelf", ".shelf"), ("sky", ".sky"),
                       ("chart", ".chartbox"), ("chat", ".chat"),
                       ("devices", ".devices"), ("matrix", ".cmp"),
                       ("tunisia", ".tn"), ("pipe", ".pipe")]:
            if stage.select_one(sel):
                kind = k
                break

    cards = []
    for c in sec.select(".icard"):
        cards.append((txt(c.select_one(".k")), txt(c.select_one("h3")), txt(c.select_one("p"))))

    info = sec.select_one(".info")
    cols = 3
    if info:
        cl = info.get("class", [])
        cols = 4 if "c4" in cl else (2 if "c2" in cl else 3)

    SLIDES.append({
        "eyebrow": txt(sec.select_one(".eyebrow")),
        "line": sec.get("data-line", ""),
        "sub": sec.get("data-sub", ""),
        "kind": kind,
        "cards": cards,
        "cols": cols,
        "dense": bool(info and "dense" in info.get("class", [])),
        "chips": [txt(c) for c in sec.select(".chip")],
        "sec": sec,
    })

print("parsed %d slides from html" % len(SLIDES))

# =====================================================================
# 2 · PPTX SCAFFOLD
# =====================================================================
prs = Presentation()
prs.slide_width, prs.slide_height = px(SW), px(SH)
BLANK = prs.slide_layouts[6]


def rect(sl, x, y, w, h, fill=None, line=None, lw=1.0, radius=0.0, shape=None):
    s = sl.shapes.add_shape(
        shape or (MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE),
        px(x), px(y), px(w), px(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    s.shadow.inherit = False
    if radius:
        try: s.adjustments[0] = radius
        except Exception: pass
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    return s


def oval(sl, x, y, d, fill=None, line=None, lw=1.0):
    return rect(sl, x, y, d, d, fill, line, lw, shape=MSO_SHAPE.OVAL)


def label(shape, s, size, color, bold=False, font=MONO, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = font
    return shape


def tbox(sl, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0, after=0):
    """paras: list of paragraphs; each = list of (text, size, color, bold, font)."""
    b = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if not isinstance(paras[0], list):
        paras = [paras]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        if after: p.space_after = Pt(after)
        for (s, size, color, bold, font) in para:
            r = p.add_run(); r.text = s
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font
    return b


# ---------------------------------------------------------------- pill
def pill(sl, x, y, w=56, h=23):
    """The two-tone capsule used all over the html deck."""
    g = rect(sl, x, y, w / 2.0, h, SIGNAL, None, radius=0.5)
    d = rect(sl, x + w / 2.0, y, w / 2.0, h, PILL_W, None, radius=0.5)
    return [g, d]


# =====================================================================
# 3 · STAGE VISUALS  (redrawn from the html scenes)
# =====================================================================
def v_radar(sl, cx, cy, size):
    d = size
    oval(sl, cx - d / 2, cy - d / 2, d, None, LINE)
    oval(sl, cx - d * .32, cy - d * .32, d * .64, None, LINE)
    oval(sl, cx - d * .14, cy - d * .14, d * .28, None, LINE)
    rect(sl, cx - d / 2, cy, d, 1, LINE)
    rect(sl, cx, cy - d / 2, 1, d, LINE)
    # sweep arm
    arm = rect(sl, cx, cy - 1, d / 2, 2, SIGNAL)
    arm.rotation = -35
    for (ox, oy, c) in [(-.22, -.20, SIGNAL2), (.24, -.30, SIGNAL2), (.16, .26, SIGNAL2)]:
        oval(sl, cx + d * ox, cy + d * oy, 10, c)
    b = sl.shapes.add_textbox(px(cx - d / 2), px(cy - 22), px(d), px(44))
    tf = b.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for s, col in [("Sentinelle", TEXT), ("Rx", SIGNAL)]:
        r = p.add_run(); r.text = s
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = col
        r.font.name = "Trebuchet MS"
    return b


def v_shelf(sl, cx, cy, h):
    w = 440
    x0 = cx - w / 2
    rows = 2
    gap = 20
    rowh = 42
    top = cy - (rows * rowh + (rows - 1) * gap) / 2
    for r in range(rows):
        y = top + r * (rowh + gap)
        n = 6 if r == 0 else 5
        for i in range(n):
            pill(sl, x0 + 14 + i * 70, y + rowh - 26)
        rect(sl, x0, y + rowh, w, 3, LINE)
    st = rect(sl, cx - 105, cy - 26, 210, 52, None, CRIT, 3, radius=0.12)
    label(st, "RUPTURE", 22, CRIT, True)
    st.rotation = -9
    return st


def v_sky(sl, cx, cy, h):
    top = cy - h / 2
    for (dx, dy, d) in [(-46, 24, 64), (-2, 0, 86), (54, 18, 70)]:
        oval(sl, cx + dx, top + dy, d, FOG)
    rect(sl, cx - 74, top + 42, 152, 52, FOG, radius=0.5)
    bolt = rect(sl, cx - 46, top + 96, 16, 34, WARN, shape=MSO_SHAPE.LIGHTNING_BOLT)
    for i, ox in enumerate([-70, -24, 22, 66]):
        p = pill(sl, cx + ox, top + 118 + (i % 2) * 26, 44, 18)
        for s in p: s.rotation = 90
    rect(sl, cx - 100, top + h - 14, 200, 12, PANEL2, radius=0.5)
    return bolt


def v_pipe(sl, cx, cy, nodes):
    n = len(nodes)
    bulb, gap = 62, 34
    total = n * bulb + (n - 1) * gap
    x0 = cx - total / 2
    last = None
    for i, (glyph, lab) in enumerate(nodes):
        x = x0 + i * (bulb + gap)
        b = oval(sl, x, cy - bulb / 2, bulb, None, LINE, 1.5)
        label(b, glyph, 20, SIGNAL2)
        lb = sl.shapes.add_textbox(px(x - 18), px(cy + bulb / 2 + 8), px(bulb + 36), px(18))
        tf = lb.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = lab.upper()
        r.font.size = Pt(8.5); r.font.color.rgb = FOG; r.font.name = MONO
        if i < n - 1:
            rect(sl, x + bulb, cy - 1, gap, 2, LINE)
            oval(sl, x + bulb + gap / 2 - 5, cy - 5, 10, SIGNAL2)
        last = b
    return last


def v_chart(sl, cx, cy, w, h):
    x0, y0 = cx - w / 2, cy - h / 2
    for g in range(1, 5):
        rect(sl, x0, y0 + h * g / 5.0, w, 1, LINE)
    split = w * 0.62
    rect(sl, x0 + split, y0, 1, h, FOG)
    tbox(sl, x0 + split + 8, y0 + 2, 90, 16, [("TODAY", 8.5, FOG, False, MONO)])

    import math
    def yv(t):
        return y0 + h * .70 - math.sin(t * 3.1) * h * .10 - t * h * .30
    pts = [(x0 + (i / 40.0) * w, yv(i / 40.0)) for i in range(41)]
    hist = [p for p in pts if p[0] <= x0 + split]
    fut = [p for p in pts if p[0] >= x0 + split]

    # confidence cone
    fb = sl.shapes.build_freeform(px(fut[0][0]), px(fut[0][1]))
    fb.add_line_segments([(px(x), px(y - 6 - (i * 1.6))) for i, (x, y) in enumerate(fut[1:])], close=False)
    fb.add_line_segments([(px(x), px(y + 6 + (i * 1.6)))
                          for i, (x, y) in reversed(list(enumerate(fut[1:])))], close=True)
    cone = fb.convert_to_shape()
    cone.fill.solid(); cone.fill.fore_color.rgb = PANEL2
    cone.line.fill.background(); cone.shadow.inherit = False

    def line_shape(seq, dash=False):
        b = sl.shapes.build_freeform(px(seq[0][0]), px(seq[0][1]))
        b.add_line_segments([(px(x), px(y)) for x, y in seq[1:]], close=False)
        s = b.convert_to_shape()
        s.fill.background(); s.line.color.rgb = SIGNAL; s.line.width = Pt(2.2)
        s.shadow.inherit = False
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        return s

    line_shape(hist)
    line_shape(fut, dash=True)
    oval(sl, pts[-1][0] - 5, pts[-1][1] - 5, 10, SIGNAL)

    badge = rect(sl, x0 + w - 232, y0 + 6, 226, 26, PANEL, SIGNAL, radius=0.12)
    label(badge, "Insuline : 2 000 -> 2 700 / semaine", 9.5, SIGNAL)
    return badge


def v_chat(sl, cx, cy, h, fr, ar):
    w = 620
    x0 = cx - w / 2
    top = cy - h / 2
    b1 = rect(sl, x0, top, w * .82, 52, PANEL2, LINE, radius=0.22)
    tf = b1.text_frame; tf.margin_left = tf.margin_right = px(16)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = fr
    r.font.size = Pt(12.5); r.font.color.rgb = TEXT; r.font.name = "Georgia"
    tbox(sl, x0 + 6, top - 15, 200, 14, [("POURQUOI ?", 8, SIGNAL, True, MONO)])

    b2 = rect(sl, x0 + w * .18, top + 66, w * .82, 46, PANEL2, LINE, radius=0.22)
    tf = b2.text_frame; tf.margin_left = tf.margin_right = px(16)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = ar
    r.font.size = Pt(13); r.font.color.rgb = TEXT_DIM; r.font.name = SANS
    tbox(sl, x0 + w - 206, top + 51, 200, 14, [("لماذا؟", 8, SIGNAL, True, MONO)],
         align=PP_ALIGN.RIGHT)

    for i, bw in enumerate([100, 70, 50, 32, 18]):
        rect(sl, cx - 135 + sum([100, 70, 50, 32, 18][:i]) + i * 7, top + 124, bw, 7,
             SIGNAL, radius=0.5)
    return b2


def v_devices(sl, cx, cy, h):
    specs = [("Command center", 280, 176), ("Pharmacy", 196, 142), ("Citizen", 104, 178)]
    gap = 32
    total = sum(s[1] for s in specs) + gap * 2
    x = cx - total / 2
    last = None
    for name, w, hh in specs:
        y = cy - hh / 2
        f = rect(sl, x, y, w, hh, INK, LINE, 1.5, radius=0.06)
        inner = rect(sl, x + 9, y + 9, w - 18, hh - 18, PANEL2, radius=0.05)
        if name == "Command center":
            rect(sl, x + 18, y + 18, (w - 36) * .6, 7, LINE, radius=0.5)
            scr = rect(sl, x + 18, y + 34, w - 36, hh - 54, PANEL, radius=0.04)
            for (ox, oy, c) in [(.18, .26, OK), (.44, .16, WARN), (.60, .52, CRIT),
                                (.30, .64, OK), (.74, .34, HOT)]:
                oval(sl, x + 18 + (w - 36) * ox, y + 34 + (hh - 54) * oy, 9, c)
        elif name == "Pharmacy":
            rect(sl, x + 16, y + 16, (w - 32) * .85, 7, LINE, radius=0.5)
            bw = (w - 32 - 5 * 4) / 6.0
            for i in range(6):
                bh = [18, 30, 24, 38, 28, 44][i]
                rect(sl, x + 16 + i * (bw + 4), y + hh - 30 - bh, bw, bh, SIGNAL, radius=0.1)
        else:
            rect(sl, x + 14, y + 14, (w - 28) * .45, 6, LINE, radius=0.5)
            scr = rect(sl, x + 14, y + 26, w - 28, hh - 56, PANEL, radius=0.05)
            pin = rect(sl, x + w / 2 - 7, y + hh / 2 - 14, 14, 14, SIGNAL, shape=MSO_SHAPE.TEAR)
            pin.rotation = 225
        tbox(sl, x - 20, y + hh + 8, w + 40, 16, [(name.upper(), 8.5, FOG, False, MONO)],
             align=PP_ALIGN.CENTER)
        last = f
        x += w + gap
    return last


def v_matrix(sl, cx, cy, h, sec):
    cells = [txt(c) for c in sec.select(".cmp > div")]
    heads = cells[1:5]
    body = cells[5:]
    labw, cw, rh = 210, 118, 30
    tot = labw + 4 * cw
    x0 = cx - tot / 2
    y0 = cy - h / 2
    for i, hd in enumerate(heads):
        tbox(sl, x0 + labw + i * cw, y0, cw, 16,
             [(hd.upper(), 8, SIGNAL if i == 3 else FOG, i == 3, MONO)], align=PP_ALIGN.CENTER)
    last = None
    for r in range(4):
        row = body[r * 5:(r + 1) * 5]
        if len(row) < 5: break
        y = y0 + 22 + r * (rh + 6)
        tbox(sl, x0, y + 7, labw - 12, 18, [(row[0], 10.5, TEXT_DIM, False, SANS)],
             align=PP_ALIGN.RIGHT)
        for c in range(4):
            us = c == 3
            cell = rect(sl, x0 + labw + c * cw + 6, y, cw - 12, rh,
                        PANEL2 if us else None, SIGNAL if us else LINE, radius=0.14)
            label(cell, row[c + 1], 11, SIGNAL2 if us else FOG, us)
            last = cell
    return last


def v_tunisia(sl, cx, cy, h):
    pts = [(150,12),(196,26),(208,62),(184,96),(196,128),(176,164),(188,210),
           (170,258),(178,300),(156,352),(138,368),(124,330),(132,282),(116,236),
           (128,190),(112,148),(124,104),(108,66),(126,30)]
    sx, sy = h / 380.0, h / 380.0
    ox, oy = cx - 150 * sx, cy - h / 2
    P = [(ox + a * sx, oy + b * sy) for a, b in pts]
    b = sl.shapes.build_freeform(px(P[0][0]), px(P[0][1]))
    b.add_line_segments([(px(a), px(bb)) for a, bb in P[1:]], close=True)
    shp = b.convert_to_shape()
    shp.fill.solid(); shp.fill.fore_color.rgb = PANEL2
    shp.line.color.rgb = SIGNAL; shp.line.width = Pt(2)
    shp.shadow.inherit = False

    for (fx, fy, name, col) in [(.56, .16, "Tunis", OK), (.57, .46, "Sfax", HOT),
                                (.44, .66, "Gabès", CRIT)]:
        x, y = ox + 300 * sx * fx, oy + h * fy
        oval(sl, x, y, 13, col)
        tbox(sl, x + 18, y - 3, 120, 16, [(name, 9.5, TEXT_DIM, False, MONO)])
    pin = rect(sl, ox + 300 * sx * .52, oy + h * .07, 20, 20, SIGNAL, shape=MSO_SHAPE.TEAR)
    pin.rotation = 225
    return pin


# =====================================================================
# 4 · ANIMATION  (raw DrawingML timing + morph transitions)
# =====================================================================
NS = ('xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
      'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
      'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
      'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"')


def _frag(xml):
    return etree.fromstring(xml)


def add_morph(slide, ms=900):
    """Morph transition, with a fade fallback for older PowerPoint."""
    xml = (
        '<mc:AlternateContent %s>'
        '<mc:Choice Requires="p14">'
        '<p:transition spd="slow" p14:dur="%d"><p14:morph option="byObject"/></p:transition>'
        '</mc:Choice>'
        '<mc:Fallback>'
        '<p:transition spd="slow"><p:fade/></p:transition>'
        '</mc:Fallback>'
        '</mc:AlternateContent>' % (NS, ms)
    )
    slide._element.append(_frag(xml))


def _effect(cid, spid, delay, dur=500, letters=None, after=True):
    """One entrance effect. letters=ms -> reveal letter by letter (robot typing)."""
    it = ('<p:iterate type="lt"><p:tmAbs val="%d"/></p:iterate>' % letters) if letters else ""
    return (
      '<p:par><p:cTn id="%d" fill="hold">'
      '<p:stCondLst><p:cond delay="%d"/></p:stCondLst><p:childTnLst>'
        '<p:par><p:cTn id="%d" fill="hold">'
        '<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
          '<p:par><p:cTn id="%d" presetID="10" presetClass="entr" presetSubtype="0" '
                'fill="hold" grpId="0" nodeType="%s">'
          '<p:stCondLst><p:cond delay="0"/></p:stCondLst>%s<p:childTnLst>'
            '<p:set><p:cBhvr>'
              '<p:cTn id="%d" dur="1" fill="hold">'
              '<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
              '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl>'
              '<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            '</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
            '<p:animEffect transition="in" filter="fade"><p:cBhvr>'
              '<p:cTn id="%d" dur="%d"/>'
              '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl>'
            '</p:cBhvr></p:animEffect>'
          '</p:childTnLst></p:cTn></p:par>'
        '</p:childTnLst></p:cTn></p:par>'
      '</p:childTnLst></p:cTn></p:par>'
      % (cid, delay, cid + 1, cid + 2, "afterEffect" if after else "withEffect",
         it, cid + 3, spid, cid + 4, dur, spid)
    )


def animate(slide, steps):
    """steps = [(shape, delay_ms, dur_ms, letters_ms_or_None), ...] in order."""
    steps = [s for s in steps if s[0] is not None]
    if not steps:
        return
    body, cid = [], 10
    for shape, delay, dur, letters in steps:
        body.append(_effect(cid, shape.shape_id, delay, dur, letters))
        cid += 10
    xml = (
      '<p:timing %s><p:tnLst>'
      '<p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
      '<p:childTnLst>'
        '<p:seq concurrent="1" nextAc="seek">'
        '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>%s</p:childTnLst>'
        '<p:prevCondLst><p:cond evt="onPrev" delay="0">'
          '<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        '<p:nextCondLst><p:cond evt="onNext" delay="0">'
          '<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        '</p:cTn></p:seq>'
      '</p:childTnLst></p:cTn></p:par>'
      '</p:tnLst></p:timing>' % (NS, "".join(body))
    )
    slide._element.append(_frag(xml))


# =====================================================================
# 5 · BUILD EACH SLIDE  (same vertical rhythm as the html)
# =====================================================================
PAD_X, PAD_T = 54, 34
CONTENT_W = SW - PAD_X * 2

def robot(sl, x, y):
    """The SRX-01 narrator, rebuilt from the html robot."""
    rect(sl, x + 42, y, 3, 14, FOG)
    oval(sl, x + 38, y - 8, 11, SIGNAL)
    head = rect(sl, x + 6, y + 14, 76, 58, PANEL2, LINE, 1.5, radius=0.24)
    oval(sl, x + 26, y + 34, 12, SIGNAL2)
    oval(sl, x + 52, y + 34, 12, SIGNAL2)
    for i in range(5):
        rect(sl, x + 26 + i * 8, y + 54, 5, 4, FOG, radius=0.5)
    rect(sl, x + 22, y + 72, 44, 12, PANEL, LINE, 1.5, radius=0.3)
    tbox(sl, x - 6, y + 88, 100, 14, [("SRX-01", 7.5, FOG, True, MONO)], align=PP_ALIGN.CENTER)
    return head


def speech_runs(line):
    out = []
    for i, part in enumerate(line.split("|")):
        if part:
            out.append((part, 19, SIGNAL2 if i % 2 else TEXT, True, MONO))
    return out


for idx, S in enumerate(SLIDES):
    sl = prs.slides.add_slide(BLANK)

    # frame
    rect(sl, 0, 0, SW, SH, INK)
    rect(sl, 8, 6, SW - 16, SH - 12, PANEL, LINE, radius=0.02)

    # eyebrow
    eb = tbox(sl, PAD_X, PAD_T, CONTENT_W, 18, [(S["eyebrow"].upper(), 9.5, SIGNAL, True, MONO)])
    rect(sl, PAD_X, PAD_T + 22, CONTENT_W, 1, LINE)

    n_cards = len(S["cards"])
    cols = S["cols"]
    rows = (n_cards + cols - 1) // cols if n_cards else 0

    # vertical rhythm mirrors the html: stage flexes, cards pinned near the base
    if rows == 0:
        # no cards: the chip row (if any) takes the card band instead
        card_h = 0
        cards_top = 610 if S["chips"] else 664
    elif rows == 1:
        card_h = 128; cards_top = 664 - card_h
    else:
        card_h = 112; cards_top = 664 - (rows * card_h + (rows - 1) * 10)

    sub_top = cards_top - 34
    voice_top = sub_top - 84
    stage_top, stage_bot = PAD_T + 34, voice_top - 10
    scy = (stage_top + stage_bot) / 2.0
    sh_avail = stage_bot - stage_top

    # ---- stage visual ----
    k = S["kind"]
    hero = None
    if k == "radar":
        hero = v_radar(sl, SW / 2, scy, min(sh_avail, 218))
    elif k == "shelf":
        hero = v_shelf(sl, SW / 2, scy, sh_avail)
    elif k == "sky":
        hero = v_sky(sl, SW / 2, scy, min(sh_avail, 200))
    elif k == "chart":
        hero = v_chart(sl, SW / 2, scy, 760, min(sh_avail, 190))
    elif k == "chat":
        fr = S["sec"].select_one(".bub.fr")["data-type"]
        ar = S["sec"].select_one(".bub.ar")["data-type"]
        hero = v_chat(sl, SW / 2, scy, min(sh_avail, 180), fr, ar)
    elif k == "devices":
        hero = v_devices(sl, SW / 2, scy, sh_avail)
    elif k == "matrix":
        hero = v_matrix(sl, SW / 2, scy, min(sh_avail, 180), S["sec"])
    elif k == "tunisia":
        hero = v_tunisia(sl, SW / 2, scy, min(sh_avail, 200))
    elif k == "pipe":
        nodes = [(txt(n.select_one(".bulb")), txt(n.select_one(".lbl")))
                 for n in S["sec"].select(".pipe .node")]
        hero = v_pipe(sl, SW / 2, scy, nodes)

    # ---- robot + speech bubble ----
    rb = robot(sl, PAD_X, voice_top - 4)
    bub = rect(sl, PAD_X + 104, voice_top, CONTENT_W - 104, 74, PANEL2, LINE, radius=0.16)
    tf = bub.text_frame
    tf.margin_left = tf.margin_right = px(20)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.line_spacing = 1.2
    for (s, size, color, bold, font) in speech_runs(S["line"]):
        r = p.add_run(); r.text = s
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font

    sub = tbox(sl, PAD_X + 104, sub_top, CONTENT_W - 104, 20,
               [(S["sub"], 10, FOG, False, MONO)])

    # ---- cards ----
    card_shapes = []
    if n_cards:
        gap = 10
        cw = (CONTENT_W - gap * (cols - 1)) / float(cols)
        for i, (kick, head, body) in enumerate(S["cards"]):
            c, r = i % cols, i // cols
            x = PAD_X + c * (cw + gap)
            y = cards_top + r * (card_h + gap)
            box_ = rect(sl, x, y, cw, card_h, PANEL2, LINE, radius=0.08)
            tbox(sl, x + 16, y + 12, cw - 32, card_h - 20,
                 [[(kick.upper(), 8.5, SIGNAL, True, MONO)],
                  [(head, 12, TEXT, True, SANS)],
                  [(body, 10, TEXT_DIM, False, SANS)]],
                 spacing=1.14, after=3)
            card_shapes.append(box_)

    # ---- chips (ask slide) ----
    chip_shapes = []
    if S["chips"]:
        gap = 12
        cwid = (CONTENT_W - gap * (len(S["chips"]) - 1)) / float(len(S["chips"]))
        for i, c in enumerate(S["chips"]):
            ch = rect(sl, PAD_X + i * (cwid + gap), cards_top - 10, cwid, 44,
                      PANEL2, SIGNAL, radius=0.4)
            label(ch, c, 11, SIGNAL2, True)
            chip_shapes.append(ch)

    # ---- footer ----
    tbox(sl, PAD_X, SH - 34, 500, 16,
         [("SENTINELLERX  ·  LA MÉTÉO DES MÉDICAMENTS", 8.5, FOG, False, MONO)])
    tbox(sl, SW - PAD_X - 200, SH - 34, 200, 16,
         [("%02d / %02d" % (idx + 1, len(SLIDES)), 8.5, FOG, False, MONO)],
         align=PP_ALIGN.RIGHT)

    # ---- animation: hero -> robot -> typed speech -> sub -> cards ----
    steps = [(hero, 200, 600, None),
             (rb, 0, 350, None),
             (bub, 0, 900, 55),          # 55ms per letter == robot typing
             (sub, 200, 400, None)]
    for cs in card_shapes:
        steps.append((cs, 120, 380, None))
    for cs in chip_shapes:
        steps.append((cs, 150, 380, None))
    animate(sl, steps)
    add_morph(sl)

    # ---- speaker notes: the spoken line, plain ----
    spoken = S["line"].replace("|", "")
    spoken = spoken[0] + spoken[1:].lower() if spoken else ""
    sl.notes_slide.notes_text_frame.text = spoken + ("\n\n" + S["sub"] if S["sub"] else "")

prs.save(OUT)
print("saved:", OUT)
